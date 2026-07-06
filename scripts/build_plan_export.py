import argparse
import json
from pathlib import Path
from xml.sax.saxutils import escape

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PURPLE = colors.HexColor("#2f1c3d")
GOLD = colors.HexColor("#fabe21")
BLUE = colors.HexColor("#88c7d5")
DARK_GRAY = colors.HexColor("#3f454a")
LIGHT_BG = colors.HexColor("#f7fbfc")
SOFT_BG = colors.HexColor("#fbf8fd")
LOGO_PATH = Path(__file__).resolve().parents[1] / "www" / "baltimore-city-logo.png"


def as_dict(value):
    return value if isinstance(value, dict) else {}


def as_list(value):
    return value if isinstance(value, list) else []


def clean(value):
    if value is None:
        return ""
    return escape(str(value))


def raw(value):
    if value is None:
        return ""
    return str(value)


def fy_label(value):
    try:
        return f"FY{int(value) % 100:02d}"
    except (TypeError, ValueError):
        return "FY"


def add_pdf_header(canvas, doc):
    canvas.saveState()
    canvas.setFillColor(PURPLE)
    canvas.rect(0, letter[1] - 0.42 * inch, letter[0], 0.42 * inch, stroke=0, fill=1)
    canvas.setFillColor(GOLD)
    canvas.rect(0, letter[1] - 0.48 * inch, letter[0], 0.06 * inch, stroke=0, fill=1)
    canvas.setFont("Helvetica", 8)
    canvas.setFillColor(colors.white)
    if LOGO_PATH.exists():
        canvas.drawImage(str(LOGO_PATH), 0.5 * inch, letter[1] - 0.37 * inch, width=0.25 * inch, height=0.25 * inch, mask="auto")
        text_x = 0.82 * inch
    else:
        text_x = 0.6 * inch
    canvas.drawString(text_x, letter[1] - 0.27 * inch, "Beacon | Baltimore City Performance & Budgeting")
    canvas.setFillColor(DARK_GRAY)
    canvas.drawRightString(letter[0] - 0.6 * inch, 0.35 * inch, f"Page {doc.page}")
    canvas.restoreState()


def pdf_title_block(payload, styles):
    title = Paragraph(f"{clean(payload.get('agency_name', 'Agency'))} Performance Plan", styles["Title"])
    subtitle = Paragraph("Beacon | Baltimore City Performance & Budgeting", styles["Subtitle"])
    text_block = [title, subtitle]
    if LOGO_PATH.exists():
        logo = Image(str(LOGO_PATH), width=0.58 * inch, height=0.58 * inch)
        table = Table([[logo, text_block]], colWidths=[0.75 * inch, 6.05 * inch], hAlign="LEFT")
        table.setStyle(TableStyle([
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 0),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ]))
        return [table]
    return [title, subtitle]


def section_band(label, styles):
    table = Table([[Paragraph(clean(label), styles["SectionBand"])]], colWidths=[6.8 * inch], hAlign="LEFT")
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), PURPLE),
        ("TEXTCOLOR", (0, 0), (-1, -1), colors.white),
        ("LEFTPADDING", (0, 0), (-1, -1), 8),
        ("RIGHTPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return [Spacer(1, 0.08 * inch), table, Spacer(1, 0.08 * inch)]


def measure_status_label(measure):
    status = raw(measure.get("validation_status") or measure.get("approval_status") or "")
    return status if status else "Not Validated"


def measure_meta_line(measure):
    return " | ".join([part for part in [
        raw(measure.get("type")),
        raw(measure.get("direction")),
        measure_status_label(measure),
    ] if part])


def pdf_measure_table(measure, styles):
    columns = as_list(measure.get("columns"))
    values = as_list(measure.get("values"))
    if len(values) < len(columns):
        values = values + [""] * (len(columns) - len(values))
    if not columns:
        columns = ["Measure"]
        values = [""]
    data = [[Paragraph(clean(col), styles["TableHeader"]) for col in columns]]
    data.append([Paragraph(clean(value), styles["TableCell"]) for value in values[: len(columns)]])
    table = Table(data, repeatRows=1, hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), SOFT_BG),
                ("TEXTCOLOR", (0, 0), (-1, 0), DARK_GRAY),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 6.5),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#d8e2e8")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return table


def pdf_score_table(scores, styles):
    rows = as_list(scores)
    if not rows:
        return []
    data = [[
        Paragraph("Criterion", styles["TableHeader"]),
        Paragraph("Score", styles["TableHeader"]),
        Paragraph("Weighted", styles["TableHeader"]),
        Paragraph("Reviewer notes", styles["TableHeader"]),
    ]]
    for row in rows:
        row = as_dict(row)
        data.append([
            Paragraph(clean(row.get("criterion")), styles["TableCell"]),
            Paragraph(clean(row.get("score")), styles["TableCell"]),
            Paragraph(clean(row.get("weighted_score")), styles["TableCell"]),
            Paragraph(clean(row.get("notes")), styles["TableCell"]),
        ])
    table = Table(data, colWidths=[2.25 * inch, 0.6 * inch, 0.72 * inch, 3.0 * inch], repeatRows=1, hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), SOFT_BG),
                ("TEXTCOLOR", (0, 0), (-1, 0), DARK_GRAY),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 6.5),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#d8e2e8")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return [Paragraph("REVIEW SCORES", styles["Eyebrow"]), table, Spacer(1, 0.1 * inch)]


def pdf_review_summary(review, styles):
    notes = as_list(review.get("notes"))
    rows = [
        ("Reviewer", review.get("reviewer")),
        ("Overall score", review.get("score")),
    ]
    data = [[Paragraph(clean(label), styles["MetaLabel"]), Paragraph(clean(value), styles["MetaValue"])] for label, value in rows if value]
    if notes:
        data.append([Paragraph("Feedback", styles["MetaLabel"]), Paragraph("<br/>".join(f"- {clean(note)}" for note in notes), styles["MetaValue"])])
    if not data:
        return []
    table = Table(data, colWidths=[1.15 * inch, 5.65 * inch], hAlign="LEFT")
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), SOFT_BG),
        ("BOX", (0, 0), (-1, -1), 0.45, colors.HexColor("#ded7e7")),
        ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#ebe4f1")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return [Paragraph("REVIEW SUMMARY", styles["Eyebrow"]), table, Spacer(1, 0.12 * inch)]


def pdf_meta_table(payload, review, styles):
    rows = [
        ("Status", payload.get("status")),
        ("Version", payload.get("version")),
        ("Plan contact", payload.get("agency_contact")),
    ]
    data = []
    for label, value in rows:
        if value is None or str(value).strip() == "":
            continue
        data.append([
            Paragraph(clean(label), styles["MetaLabel"]),
            Paragraph(clean(value), styles["MetaValue"]),
        ])
    if not data:
        return []
    table = Table(data, colWidths=[1.15 * inch, 5.65 * inch], hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), LIGHT_BG),
                ("BOX", (0, 0), (-1, -1), 0.45, colors.HexColor("#d8e2e8")),
                ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#e4ebef")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
        )
    )
    return [table, Spacer(1, 0.12 * inch)]


def build_pdf(payload, output):
    overview = as_dict(payload.get("overview"))
    review = as_dict(payload.get("review"))
    doc = SimpleDocTemplate(
        output,
        pagesize=letter,
        rightMargin=0.55 * inch,
        leftMargin=0.55 * inch,
        topMargin=0.78 * inch,
        bottomMargin=0.65 * inch,
    )
    base = getSampleStyleSheet()
    styles = {
        "Title": ParagraphStyle("Title", parent=base["Title"], textColor=PURPLE, fontSize=22, leading=26, spaceAfter=5, alignment=0),
        "Subtitle": ParagraphStyle("Subtitle", parent=base["Normal"], textColor=DARK_GRAY, fontSize=9, leading=12, spaceAfter=10),
        "H1": ParagraphStyle("H1", parent=base["Heading1"], textColor=PURPLE, fontSize=16, leading=20, spaceBefore=12, spaceAfter=8),
        "H2": ParagraphStyle("H2", parent=base["Heading2"], textColor=PURPLE, fontSize=12, leading=15, spaceBefore=8, spaceAfter=5),
        "SectionBand": ParagraphStyle("SectionBand", parent=base["BodyText"], textColor=colors.white, fontName="Helvetica-Bold", fontSize=10, leading=12, alignment=0),
        "Eyebrow": ParagraphStyle("Eyebrow", parent=base["BodyText"], textColor=DARK_GRAY, fontName="Helvetica-Bold", fontSize=7.5, leading=9, spaceBefore=4, spaceAfter=4),
        "MeasureTitle": ParagraphStyle("MeasureTitle", parent=base["BodyText"], textColor=colors.HexColor("#1f1f23"), fontName="Helvetica-Bold", fontSize=9, leading=11, spaceBefore=4, spaceAfter=2),
        "Body": ParagraphStyle("Body", parent=base["BodyText"], fontSize=9.5, leading=13, spaceAfter=6),
        "Small": ParagraphStyle("Small", parent=base["BodyText"], textColor=DARK_GRAY, fontSize=8, leading=11, spaceAfter=4),
        "MetaLabel": ParagraphStyle("MetaLabel", parent=base["BodyText"], textColor=DARK_GRAY, fontName="Helvetica-Bold", fontSize=7.5, leading=10),
        "MetaValue": ParagraphStyle("MetaValue", parent=base["BodyText"], textColor=PURPLE, fontName="Helvetica-Bold", fontSize=8.5, leading=11),
        "TableHeader": ParagraphStyle("TableHeader", parent=base["BodyText"], textColor=DARK_GRAY, fontSize=6.5, leading=8),
        "TableCell": ParagraphStyle("TableCell", parent=base["BodyText"], fontSize=6.5, leading=8),
    }
    story = []
    story.extend(pdf_title_block(payload, styles))
    story.extend(pdf_meta_table(payload, review, styles))
    if payload.get("include_review", True):
        story.extend(pdf_review_summary(review, styles))

    story.extend(section_band("Overview & Vision", styles))
    story.append(Paragraph(f"<b>Overview:</b> {clean(overview.get('overview'))}", styles["Body"]))
    story.append(Paragraph(f"<b>Vision:</b> {clean(overview.get('vision'))}", styles["Body"]))
    if overview.get("web_address"):
        story.append(Paragraph(f"<b>Web address:</b> {clean(overview.get('web_address'))}", styles["Body"]))
    if payload.get("include_review", True):
        story.extend(pdf_score_table(payload.get("overview_scores"), styles))

    story.extend(section_band("Agency Goals", styles))
    for goal in as_list(payload.get("goals")):
        goal = as_dict(goal)
        story.append(Paragraph(clean(goal.get("title")), styles["H2"]))
        story.append(Spacer(1, 0.04 * inch))
        if goal.get("initiatives"):
            story.append(Paragraph("INITIATIVES", styles["Eyebrow"]))
            for initiative in as_list(goal.get("initiatives")):
                story.append(Paragraph(f"- {clean(initiative)}", styles["Body"]))
            story.append(Spacer(1, 0.08 * inch))
        if goal.get("kpis"):
            story.append(Paragraph("KEY PERFORMANCE INDICATORS", styles["Eyebrow"]))
            story.append(Spacer(1, 0.03 * inch))
            for measure in as_list(goal.get("kpis")):
                measure = as_dict(measure)
                story.append(Paragraph(clean(measure.get("title")), styles["MeasureTitle"]))
                story.append(Paragraph(clean(measure_meta_line(measure)), styles["Small"]))
                story.append(pdf_measure_table(measure, styles))
                story.append(Spacer(1, 0.08 * inch))
        if goal.get("alignment"):
            story.append(Paragraph(f"<b>Action Plan alignment:</b> {clean(goal['alignment'])}", styles["Small"]))
        if payload.get("include_review", True):
            story.extend(pdf_score_table(goal.get("review_scores"), styles))

    story.extend(section_band("Services", styles))
    for service in as_list(payload.get("services")):
        service = as_dict(service)
        story.append(Paragraph(clean(service.get("name")), styles["H2"]))
        if service.get("scoring_exempt"):
            story.append(Paragraph("Administration service: not scored and no service metrics required this cycle.", styles["Small"]))
        story.append(Paragraph(clean(service.get("description")), styles["Body"]))
        if service.get("metrics"):
            story.append(Paragraph("PERFORMANCE METRICS", styles["Eyebrow"]))
            for measure in as_list(service.get("metrics")):
                measure = as_dict(measure)
                story.append(Paragraph(clean(measure.get("title")), styles["MeasureTitle"]))
                story.append(Paragraph(clean(measure_meta_line(measure)), styles["Small"]))
                story.append(pdf_measure_table(measure, styles))
                story.append(Spacer(1, 0.08 * inch))
        if payload.get("include_review", True):
            story.extend(pdf_score_table(service.get("review_scores"), styles))

    story.extend(section_band("Risks", styles))
    risks = as_list(payload.get("risks"))
    if risks:
        for risk in risks:
            if isinstance(risk, dict):
                category = risk.get("category") or "Uncategorized"
                description = risk.get("description")
                story.append(Paragraph(f"<b>{clean(category)}:</b> {clean(description)}", styles["Body"]))
            else:
                story.append(Paragraph(f"- {clean(risk)}", styles["Body"]))
    else:
        story.append(Paragraph("- No risks available.", styles["Body"]))
    if payload.get("include_review", True):
        story.extend(pdf_score_table(payload.get("risk_scores"), styles))
        if as_list(payload.get("plan_scores")):
            story.extend(section_band("Plan-Level Review Scores", styles))
            story.extend(pdf_score_table(payload.get("plan_scores"), styles))

    doc.build(story, onFirstPage=add_pdf_header, onLaterPages=add_pdf_header)


def set_run(run, size=18, bold=False, color=PURPLE):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(color.red if hasattr(color, "red") else 47, color.green if hasattr(color, "green") else 28, color.blue if hasattr(color, "blue") else 61)


def add_textbox(slide, text, left, top, width, height, size=18, bold=False, color=(47, 28, 61)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return box


def blank_layout(prs):
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, title, Inches(0.55), Inches(0.35), Inches(8.8), Inches(0.45), 24, True)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.55), Inches(0.82), Inches(8.8), Inches(0.35), 10, False, (63, 69, 74))
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 190, 33)
    shape.line.fill.background()


def add_metric_table_to_slide(slide, measure, y):
    measure = as_dict(measure)
    columns = as_list(measure.get("columns"))
    values = as_list(measure.get("values"))
    if len(values) < len(columns):
        values = values + [""] * (len(columns) - len(values))
    if not columns:
        columns = ["Measure"]
        values = [""]
    add_textbox(slide, raw(measure.get("title")), Inches(0.7), y, Inches(8.6), Inches(0.28), 10, True, (22, 22, 22))
    add_textbox(slide, f"{raw(measure.get('type'))} | {raw(measure.get('direction'))}", Inches(0.7), y + Inches(0.25), Inches(8.6), Inches(0.25), 8, False, (63, 69, 74))
    rows, cols = 2, len(columns)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.7), y + Inches(0.58), Inches(8.6), Inches(0.72))
    table = table_shape.table
    for idx, col in enumerate(columns):
        table.cell(0, idx).text = raw(col)
        table.cell(1, idx).text = raw(values[idx])
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(6.5)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return y + Inches(1.42)


def build_pptx(payload, output, template=None):
    overview = as_dict(payload.get("overview"))
    review = as_dict(payload.get("review"))
    prs = Presentation(template) if template and Path(template).exists() else Presentation()
    layout = blank_layout(prs)

    slide = prs.slides.add_slide(layout)
    add_slide_title(slide, f"{raw(payload.get('agency_name', 'Agency'))} Performance Plan", fy_label(payload.get("fiscal_year")))
    add_textbox(slide, "Beacon | Baltimore City Performance & Budgeting", Inches(0.65), Inches(1.6), Inches(8.5), Inches(0.5), 26, True)
    add_textbox(slide, f"{raw(payload.get('status'))} | Version {raw(payload.get('version'))} | {raw(payload.get('agency_contact'))}", Inches(0.68), Inches(2.2), Inches(8.4), Inches(0.45), 13, False, (63, 69, 74))

    slide = prs.slides.add_slide(layout)
    if payload.get("include_review", True):
        add_slide_title(slide, "Overview & Reviewer Feedback", f"Overall score: {raw(review.get('score', 'Not scored'))}")
    else:
        add_slide_title(slide, "Overview")
    add_textbox(slide, "Overview", Inches(0.65), Inches(1.25), Inches(4.1), Inches(0.3), 15, True)
    add_textbox(slide, raw(overview.get("overview")), Inches(0.65), Inches(1.65), Inches(4.1), Inches(1.8), 12, False, (22, 22, 22))
    if payload.get("include_review", True):
        add_textbox(slide, "Top improvement areas", Inches(5.1), Inches(1.25), Inches(4.1), Inches(0.3), 15, True)
        add_textbox(slide, "\n".join(f"- {raw(note)}" for note in as_list(review.get("notes"))), Inches(5.1), Inches(1.65), Inches(4.2), Inches(2.1), 11, False, (22, 22, 22))

    for goal in as_list(payload.get("goals")):
        goal = as_dict(goal)
        slide = prs.slides.add_slide(layout)
        add_slide_title(slide, "Agency Goal", raw(goal.get("title"))[:95])
        y = Inches(1.25)
        if goal.get("initiatives"):
            add_textbox(slide, "Initiatives", Inches(0.65), y, Inches(8.8), Inches(0.25), 12, True, (63, 69, 74))
            add_textbox(slide, "\n".join(f"- {raw(i)}" for i in as_list(goal.get("initiatives"))), Inches(0.75), y + Inches(0.32), Inches(8.6), Inches(0.55), 10, False, (22, 22, 22))
            y += Inches(1.05)
        for measure in as_list(goal.get("kpis"))[:2]:
            y = add_metric_table_to_slide(slide, measure, y)

    slide = prs.slides.add_slide(layout)
    services = as_list(payload.get("services"))
    add_slide_title(slide, "Services & Performance Metrics", f"{len(services)} services")
    y = Inches(1.2)
    for service in services[:3]:
        service = as_dict(service)
        add_textbox(slide, raw(service.get("name")), Inches(0.65), y, Inches(8.8), Inches(0.25), 12, True)
        y += Inches(0.33)
        if service.get("scoring_exempt"):
            add_textbox(slide, "Administration service: not scored this cycle.", Inches(0.7), y, Inches(8.4), Inches(0.24), 8, False, (63, 69, 74))
            y += Inches(0.35)
        for measure in as_list(service.get("metrics"))[:1]:
            y = add_metric_table_to_slide(slide, measure, y)
        y += Inches(0.12)
        if y > Inches(6.0):
            break

    slide = prs.slides.add_slide(layout)
    risks = as_list(payload.get("risks"))
    add_slide_title(slide, "Risks", f"{len(risks)} risks identified")
    risk_lines = []
    for risk in risks:
        if isinstance(risk, dict):
            risk_lines.append(f"- {raw(risk.get('category') or 'Uncategorized')}: {raw(risk.get('description'))}")
        else:
            risk_lines.append(f"- {raw(risk)}")
    add_textbox(slide, "\n".join(risk_lines) or "No risks available.", Inches(0.75), Inches(1.3), Inches(8.5), Inches(4.8), 14, False, (22, 22, 22))

    prs.save(output)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--type", choices=["pdf", "pptx"], required=True)
    parser.add_argument("--template")
    args = parser.parse_args()
    with open(args.input, "r", encoding="utf-8") as handle:
        payload = json.load(handle)
    if args.type == "pdf":
        build_pdf(payload, args.output)
    else:
        build_pptx(payload, args.output, args.template)


if __name__ == "__main__":
    main()
